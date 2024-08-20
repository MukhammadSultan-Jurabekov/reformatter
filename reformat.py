from docx import Document
import pdfkit
from fpdf import FPDF
from PIL import Image
import pdf2image
from pptx import Presentation
import pandas as pd
import os

def word_to_pdf(docx_path, pdf_path):
    pdfkit.from_file(docx_path, pdf_path)

def txt_to_pdf(txt_path, pdf_path):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    
    with open(txt_path, 'r') as file:
        for line in file:
            pdf.cell(200, 10, txt=line.encode('latin-1', 'replace').decode('latin-1'), ln=True)
    
    pdf.output(pdf_path)

def png_to_jpeg(png_path, jpeg_path):
    image = Image.open(png_path)
    rgb_im = image.convert('RGB')
    rgb_im.save(jpeg_path, format="JPEG")

def pdf_to_pptx(pdf_path, pptx_path):
    images = pdf2image.convert_from_path(pdf_path)
    prs = Presentation()
    
    for image in images:
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        img_path = "temp_image.jpg"
        image.save(img_path, 'JPEG')
        slide.shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)
        os.remove(img_path)
    
    prs.save(pptx_path)

def excel_to_pdf(excel_path, pdf_path):
    df = pd.read_excel(excel_path)
    df.to_html("temp.html")
    pdfkit.from_file("temp.html", pdf_path)
    os.remove("temp.html")

def html_to_pdf(html_path, pdf_path):
    pdfkit.from_file(html_path, pdf_path)

def pdf_to_word(pdf_path, docx_path):
    # Placeholder: Convert PDF to Word (requires more complex logic or a third-party API)
    pass

def pdf_to_txt(pdf_path, txt_path):
    # Placeholder: Convert PDF to TXT (requires more complex logic or a third-party API)
    pass

def jpeg_to_png(jpeg_path, png_path):
    image = Image.open(jpeg_path)
    image.save(png_path, format="PNG")

def pptx_to_pdf(pptx_path, pdf_path):
    # Placeholder: Convert PPTX to PDF (requires a library like `pywin32` on Windows or a third-party service)
    pass

def main():
    # Example usage:
    word_to_pdf('example.docx', 'output.pdf')
    txt_to_pdf('example.txt', 'output.pdf')
    png_to_jpeg('example.png', 'output.jpg')
    pdf_to_pptx('example.pdf', 'output.pptx')
    excel_to_pdf('example.xlsx', 'output.pdf')
    html_to_pdf('example.html', 'output.pdf')

if __name__ == "__main__":
    main()
